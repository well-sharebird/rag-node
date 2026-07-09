"""Document parsing service with table extraction support.

Supports: PDF (with tables), DOCX (with tables), XLSX, PPTX, TXT, MD, HTML
"""
from __future__ import annotations
import io
import logging
from app.utils.exceptions import ValidationException

logger = logging.getLogger("app.services.parser")

# Extended supported formats
ALLOWED_EXTENSIONS = {
    "pdf", "docx", "xlsx", "pptx", "txt", "md", "html", "htm",
}


async def parse_document(content: bytes, format: str) -> str:
    """Parse document content into plain text. Supports tables extraction."""
    format = format.lower()
    parsers = {
        "pdf": _parse_pdf, "docx": _parse_docx,
        "xlsx": _parse_xlsx, "pptx": _parse_pptx,
        "txt": _parse_txt, "md": _parse_txt,
        "html": _parse_html, "htm": _parse_html,
    }
    parser = parsers.get(format)
    if parser is None:
        raise ValidationException(f"No parser available for format: {format}")

    try:
        text = parser(content)
        logger.debug("Parsed %s document: %d chars", format, len(text))
        return text
    except Exception as e:
        logger.exception("Failed to parse %s document: %s", format, e)
        raise ValidationException(f"Failed to parse {format} document: {e}")


# ============================================================
# PDF Parser — pdfplumber for table-aware extraction
# ============================================================

def _parse_pdf(content: bytes) -> str:
    """Parse PDF with table extraction. Falls back to PyPDF2 if pdfplumber unavailable."""
    try:
        import pdfplumber
        return _parse_pdf_with_plumber(content)
    except ImportError:
        logger.warning("pdfplumber not installed, falling back to PyPDF2 (no table support)")
        return _parse_pdf_fallback(content)


def _parse_pdf_with_plumber(content: bytes) -> str:
    """PDF parsing via pdfplumber with table detection."""
    import pdfplumber
    results = []
    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            page_texts = []

            # Extract regular text
            text = page.extract_text()
            if text:
                page_texts.append(text)

            # Extract tables
            tables = page.extract_tables()
            for table_num, table in enumerate(tables, 1):
                if table and any(any(cell for cell in row) for row in table):
                    table_text = _format_table(table, f"表格 {page_num}.{table_num}")
                    page_texts.append(table_text)

            if page_texts:
                results.append(f"## 第 {page_num} 页\n\n" + "\n\n".join(page_texts))

    return "\n\n".join(results)


def _format_table(table: list[list[str | None]], caption: str = "") -> str:
    """Format a table as readable Markdown-like text."""
    if not table:
        return ""

    lines = [f"[{caption}]"]
    # Header row
    header = table[0]
    header_str = " | ".join(str(c or "") for c in header)
    lines.append(header_str)
    lines.append("-" * len(header_str))

    # Data rows
    for row in table[1:]:
        row_str = " | ".join(str(c or "") for c in row)
        lines.append(row_str)

    return "\n".join(lines)


def _parse_pdf_fallback(content: bytes) -> str:
    """Fallback PDF parser using PyPDF2 (no table support)."""
    from PyPDF2 import PdfReader
    reader = PdfReader(io.BytesIO(content))
    texts = [page.extract_text() for page in reader.pages if page.extract_text()]
    return "\n\n".join(texts)


# ============================================================
# DOCX Parser — with table extraction
# ============================================================

def _parse_docx(content: bytes) -> str:
    """Parse DOCX with table extraction."""
    from docx import Document
    doc = Document(io.BytesIO(content))
    results = []

    for element in doc.element.body:
        tag = element.tag.split('}')[-1] if '}' in element.tag else element.tag
        if tag == 'p':
            # Paragraph
            from docx.text.paragraph import Paragraph
            para = Paragraph(element, doc)
            text = para.text.strip()
            if text:
                results.append(text)
        elif tag == 'tbl':
            # Table
            from docx.table import Table
            table = Table(element, doc)
            table_data = []
            for row in table.rows:
                table_data.append([cell.text.strip() for cell in row.cells])
            if table_data:
                results.append(_format_table(table_data, "表格"))

    # Fallback: if no elements parsed, use simple paragraph extraction
    if not results:
        results = [p.text for p in doc.paragraphs if p.text.strip()]

    return "\n\n".join(results)


# ============================================================
# XLSX Parser
# ============================================================

def _parse_xlsx(content: bytes) -> str:
    """Parse XLSX spreadsheets. Each sheet becomes a section."""
    import openpyxl

    # Try loading normally, fall back to read-only mode for complex files
    try:
        wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
    except TypeError:
        logger.warning("XLSX has complex formatting, falling back to read-only mode")
        wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True, read_only=True)

    results = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        # Filter empty rows
        filtered = [list(row) for row in rows if any(c is not None for c in row)]
        if not filtered:
            continue

        table_text = _format_table(filtered, f"工作表: {sheet_name}")
        results.append(f"## {sheet_name}\n\n{table_text}")

    wb.close()
    return "\n\n".join(results)


# ============================================================
# PPTX Parser
# ============================================================

def _parse_pptx(content: bytes) -> str:
    """Parse PPTX presentations. Each slide becomes a section."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    results = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_texts.append(text)
            if shape.has_table:
                table = shape.table
                table_data = []
                for row in table.rows:
                    table_data.append([cell.text.strip() for cell in row.cells])
                if table_data:
                    slide_texts.append(_format_table(table_data, f"幻灯片 {slide_num} 表格"))
        if slide_texts:
            results.append(f"## 幻灯片 {slide_num}\n\n" + "\n\n".join(slide_texts))

    return "\n\n".join(results)


# ============================================================
# Simple parsers
# ============================================================

def _parse_txt(content: bytes) -> str:
    return content.decode("utf-8", errors="replace")


def _parse_html(content: bytes) -> str:
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(content, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()

    # Extract tables from HTML
    for table_tag in soup.find_all("table"):
        table_data = []
        for row in table_tag.find_all("tr"):
            cells = [cell.get_text(strip=True) for cell in row.find_all(["td", "th"])]
            if cells:
                table_data.append(cells)
        if table_data:
            caption = table_tag.get("summary") or table_tag.get("caption") or "HTML 表格"
            formatted = _format_table(table_data, str(caption))
            # Replace table tag with formatted text
            # (Simple approach: insert before)
            table_tag.insert_before(formatted)

    return soup.get_text(separator="\n", strip=True)
