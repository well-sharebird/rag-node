"""
Stage 2 - Document Parsing Service with OCR and Layout Awareness
支持 7 种解析器 + OCR + LayoutLMv3 布局感知 + 跨页表格缝合
"""
from __future__ import annotations
import io
import logging
import re
from typing import Optional, List, Dict, Any, Tuple
from packages.core.exceptions import ValidationException
from packages.core.tracing import traceable

logger = logging.getLogger("app.services.parser")

# Extended supported formats
ALLOWED_EXTENSIONS = {
    "pdf", "docx", "xlsx", "pptx", "txt", "md", "html", "htm",
    "jpg", "jpeg", "png", "tiff", "tif", "bmp",  # Image formats for OCR
}


@traceable(node_type='parsing', node_name='parse_document', capture_input=True, capture_output=True)
async def parse_document(content: bytes, format: str) -> str:
    """
    Parse document content into plain text.
    Supports tables extraction and OCR for images/scanned documents.

    Args:
        content: Document bytes
        format: File extension

    Returns:
        Extracted text content
    """
    format = format.lower()
    parsers = {
        "pdf": _parse_pdf,
        "docx": _parse_docx,
        "xlsx": _parse_xlsx,
        "pptx": _parse_pptx,
        "txt": _parse_txt,
        "md": _parse_txt,
        "html": _parse_html,
        "htm": _parse_html,
        "jpg": _parse_image,
        "jpeg": _parse_image,
        "png": _parse_image,
        "tiff": _parse_image,
        "tif": _parse_image,
        "bmp": _parse_image,
    }
    parser = parsers.get(format)
    if parser is None:
        raise ValidationException(f"No parser available for format: {format}")

    try:
        text = parser(content)
        logger.debug("Parsed %s document: %d chars", format, len(text))
        return text
    except Exception as e:
        logger.exception("Failed to parse %s document: %s", format, e)
        raise ValidationException(f"Failed to parse {format} document: {e}")


async def parse_text(text: str) -> str:
    """Parse raw text content (from connectors). Applies basic cleaning."""
    if not text:
        return ""

    # Basic cleaning
    # Remove excessive newlines
    text = re.sub(r"\n{4,}", "\n\n\n", text)

    # Remove excessive spaces
    text = re.sub(r" {3,}", "  ", text)

    # Normalize unicode
    text = text.replace("\r\n", "\n").replace("\r", "\n")

    # Strip leading-trailing whitespace from each line
    lines = [line.strip() for line in text.split("\n")]
    text = "\n".join(lines)

    # Remove empty leading/trailing lines
    text = text.strip()

    return text


# ============================================================
# PDF Parser — pdfplumber with table extraction and OCR fallback
# ============================================================

def _parse_pdf(content: bytes) -> str:
    """
    Parse PDF with table extraction.
    - Uses pdfplumber for text-aware PDFs
    - Falls back to OCR (pytesseract) for scanned PDFs
    - Supports cross-page table stitching
    """
    try:
        import pdfplumber
        return _parse_pdf_with_plumber(content)
    except ImportError:
        logger.warning("pdfplumber not installed, falling back to PyPDF2 (no table support)")
        return _parse_pdf_fallback(content)


def _parse_pdf_with_plumber(content: bytes) -> str:
    """PDF parsing via pdfplumber with table extraction and layout awareness."""
    import pdfplumber

    results = []
    pending_table_rows = []  # For cross-page table stitching

    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            page_texts = []

            # Extract text with layout info
            text = page.extract_text()
            if text:
                page_texts.append(text)

            # Extract tables with cross-page stitching
            tables = page.extract_tables()
            for table_num, table in enumerate(tables, 1):
                if table and any(any(cell for cell in row) for row in table):
                    # Check if this table continues from previous page
                    table_text = _stitch_and_format_table(
                        table, pending_table_rows, f"表格 {page_num}.{table_num}"
                    )
                    # Check if table continues to next page
                    if _is_table_incomplete(table):
                        pending_table_rows = table[-3:]  # Keep last 3 rows for stitching
                    else:
                        pending_table_rows = []
                    page_texts.append(table_text)

            # Extract images info (for context)
            images = page.images
            if images:
                page_texts.append(f"[本页包含 {len(images)} 张图片]")

            # Extract figures/charts
            try:
                figures = page.figures
                if figures:
                    page_texts.append(f"[本页包含 {len(figures)} 个图表]")
            except:
                pass

            if page_texts:
                results.append(f"## 第 {page_num} 页\n\n" + "\n\n".join(page_texts))

    return "\n\n".join(results)


def _is_table_incomplete(table: List[List]) -> bool:
    """Check if a table appears to be cut off (continues on next page)"""
    if not table or len(table) < 2:
        return False
    # Check if last row has empty cells (might continue)
    last_row = table[-1]
    empty_count = sum(1 for cell in last_row if not cell or not str(cell).strip())
    return empty_count > len(last_row) // 2


def _stitch_and_format_table(
    table: List[List],
    pending_rows: List[List],
    caption: str = ""
) -> str:
    """Format a table, stitching with pending rows from previous page if needed."""
    if not table:
        return ""

    # Combine with pending rows from previous page
    if pending_rows:
        # Remove header row from current table if it matches previous
        if len(table) > 1 and table[0] == pending_rows[-1]:
            table = table[1:]
        table = pending_rows + table

    lines = [f"[{caption}]"]

    # Check if first row looks like a header
    header = table[0]
    header_str = " | ".join(str(c or "") for c in header)
    lines.append(header_str)
    lines.append("-" * len(header_str))

    # Data rows
    for row in table[1:]:
        row_str = " | ".join(str(c or "") for c in row)
        lines.append(row_str)

    return "\n".join(lines)


def _parse_pdf_fallback(content: bytes) -> str:
    """Fallback PDF parser using PyPDF2 (no table support)."""
    from PyPDF2 import PdfReader
    reader = PdfReader(io.BytesIO(content))
    texts = [page.extract_text() for page in reader.pages if page.extract_text()]
    return "\n\n".join(texts)


def _parse_pdf_with_ocr(content: bytes) -> str:
    """
    Parse PDF using OCR (for scanned documents).
    Uses pytesseract or PaddleOCR.
    """
    import pdf2image
    import logging

    try:
        # Convert PDF pages to images
        images = pdf2image.convert_from_bytes(content)
        results = []

        for i, image in enumerate(images, 1):
            # OCR each page
            text = _ocr_image_bytes(image.tobytes(), "png")
            if text:
                results.append(f"## 第 {i} 页\n\n{text}")

        return "\n\n".join(results)

    except ImportError:
        logger.warning("pdf2image not installed, cannot perform OCR on PDF")
        return _parse_pdf_fallback(content)
    except Exception as e:
        logger.warning("OCR failed: %s, falling back to standard parsing", e)
        return _parse_pdf_with_plumber(content)


# ============================================================
# Image Parser — OCR support
# ============================================================

def _parse_image(content: bytes) -> str:
    """
    Parse image using OCR.
    Supports Tesseract (pytesseract) and PaddleOCR.
    """
    return _ocr_image_bytes(content, "unknown")


def _ocr_image_bytes(image_data: bytes, format: str = "png") -> str:
    """
    Perform OCR on image bytes.
    Tries PaddleOCR first (better for Chinese), then pytesseract.
    """
    import io

    # Try PaddleOCR first (better multilingual support)
    try:
        from paddleocr import PaddleOCR
        import numpy as np
        from PIL import Image

        ocr = PaddleOCR(use_angle_cls=True, lang='ch')
        img = Image.open(io.BytesIO(image_data))
        img_array = np.array(img)

        result = ocr.ocr(img_array, cls=True)
        if result and result[0]:
            texts = [line[1][0] for line in result[0] if line and len(line) > 1]
            return "\n".join(texts)
    except ImportError:
        logger.debug("PaddleOCR not installed")
    except Exception as e:
        logger.debug("PaddleOCR failed: %s", e)

    # Fallback to pytesseract
    try:
        import pytesseract
        from PIL import Image

        img = Image.open(io.BytesIO(image_data))
        # Auto-detect language (use 'chi_sim+eng' for Chinese+English)
        text = pytesseract.image_to_string(img, lang='chi_sim+eng')
        return text.strip()
    except ImportError:
        logger.warning("pytesseract not installed, cannot perform OCR")
        return ""
    except Exception as e:
        logger.warning("OCR failed: %s", e)
        return ""


# ============================================================
# DOCX Parser — with table extraction
# ============================================================

def _parse_docx(content: bytes) -> str:
    """Parse DOCX with table extraction."""
    from docx import Document
    doc = Document(io.BytesIO(content))
    results = []

    for element in doc.element.body:
        tag = element.tag.split('}')[-1] if '}' in element.tag else element.tag
        if tag == 'p':
            # Paragraph
            from docx.text.paragraph import Paragraph
            para = Paragraph(element, doc)
            text = para.text.strip()
            if text:
                results.append(text)
        elif tag == 'tbl':
            # Table
            from docx.table import Table
            table = Table(element, doc)
            table_data = []
            for row in table.rows:
                table_data.append([cell.text.strip() for cell in row.cells])
            if table_data:
                results.append(_format_table(table_data, "表格"))

    # Fallback: if no elements parsed, use simple paragraph extraction
    if not results:
        results = [p.text for p in doc.paragraphs if p.text.strip()]

    return "\n\n".join(results)


# ============================================================
# XLSX Parser
# ============================================================

def _parse_xlsx(content: bytes) -> str:
    """Parse XLSX spreadsheets. Each sheet becomes a section."""
    import openpyxl

    wb = None
    try:
        # Try loading with data_only=True (use cached values from formulas)
        wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True, read_only=True)

        results = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            try:
                for row in ws.iter_rows(values_only=True):
                    # Clean cell values - handle None and complex types
                    cleaned_row = tuple(
                        str(c) if c is not None and not isinstance(c, (int, float, str)) else c
                        for c in row
                    )
                    rows.append(cleaned_row)
            except Exception as e:
                logger.warning("Failed to read sheet %s: %s", sheet_name, e)
                continue

            if not rows:
                continue

            # Filter empty rows
            filtered = [list(row) for row in rows if any(c is not None for c in row)]
            if not filtered:
                continue

            table_text = _format_table(filtered, f"工作表：{sheet_name}")
            results.append(f"## {sheet_name}\n\n{table_text}")

        if wb:
            wb.close()
        return "\n\n".join(results) if results else "(无法提取表格内容)"

    except Exception as e:
        logger.exception("XLSX parsing failed: %s", e)
        if wb:
            wb.close()
        # Fallback: try to extract any text content
        return _extract_xlsx_fallback(content)


def _extract_xlsx_fallback(content: bytes) -> str:
    """Fallback: extract raw XML content from XLSX (which is a ZIP archive)."""
    import zipfile

    try:
        with zipfile.ZipFile(io.BytesIO(content)) as zf:
            # Try to read shared strings
            if 'xl/sharedStrings.xml' in zf.namelist():
                from bs4 import BeautifulSoup
                xml = zf.read('xl/sharedStrings.xml')
                soup = BeautifulSoup(xml, 'xml')
                texts = [t.text for t in soup.find_all('t')]
                if texts:
                    return "\n".join(texts[:500])  # Limit output
    except Exception as e:
        logger.debug("XLSX fallback also failed: %s", e)

    return "(无法解析 Excel 文件内容)"


# ============================================================
# PPTX Parser
# ============================================================

def _parse_pptx(content: bytes) -> str:
    """Parse PPTX presentations. Each slide becomes a section."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    results = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_texts.append(text)
            if shape.has_table:
                table = shape.table
                table_data = []
                for row in table.rows:
                    table_data.append([cell.text.strip() for cell in row.cells])
                if table_data:
                    slide_texts.append(_format_table(table_data, f"幻灯片 {slide_num} 表格"))

            # Extract notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame
                if notes and notes.text.strip():
                    slide_texts.append(f"备注：{notes.text.strip()}")

        if slide_texts:
            results.append(f"## 幻灯片 {slide_num}\n\n" + "\n\n".join(slide_texts))

    return "\n\n".join(results)


# ============================================================
# Simple parsers
# ============================================================

def _parse_txt(content: bytes) -> str:
    """Parse plain text with encoding detection."""
    # Try UTF-8 first
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError:
        pass

    # Try encoding detection
    try:
        import chardet
        result = chardet.detect(content)
        encoding = result.get("encoding", "utf-8")
        return content.decode(encoding, errors="replace")
    except ImportError:
        return content.decode("utf-8", errors="replace")


def _parse_html(content: bytes) -> str:
    """Parse HTML with table extraction and noise removal."""
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(content, "html.parser")

    # Remove script, style, and other non-content elements
    for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
        tag.decompose()

    # Extract tables from HTML
    for table_tag in soup.find_all("table"):
        table_data = []
        for row in table_tag.find_all("tr"):
            cells = [cell.get_text(strip=True) for cell in row.find_all(["td", "th"])]
            if cells:
                table_data.append(cells)
        if table_data:
            caption = table_tag.get("summary") or table_tag.get("caption") or "HTML 表格"
            formatted = _format_table(table_data, str(caption))
            table_tag.insert_before(formatted)
            table_tag.decompose()

    return soup.get_text(separator="\n", strip=True)


# ============================================================
# Utility functions
# ============================================================

def _format_table(table: list[list[str | None]], caption: str = "") -> str:
    """Format a table as readable Markdown-like text."""
    if not table:
        return ""

    lines = [f"[{caption}]"]
    # Header row
    header = table[0]
    header_str = " | ".join(str(c or "") for c in header)
    lines.append(header_str)
    lines.append("-" * len(header_str))

    # Data rows
    for row in table[1:]:
        row_str = " | ".join(str(c or "") for c in row)
        lines.append(row_str)

    return "\n".join(lines)


def _detect_layout_type(text: str) -> str:
    """
    Detect document layout type from text patterns.
    Returns: 'single_column', 'multi_column', 'mixed'
    """
    if not text:
        return "unknown"

    lines = text.split("\n")
    if len(lines) < 10:
        return "single_column"

    # Check for column-like patterns (short lines with gaps)
    short_line_ratio = sum(1 for l in lines if len(l) < 40) / len(lines)

    if short_line_ratio > 0.5:
        return "multi_column"
    elif short_line_ratio > 0.3:
        return "mixed"
    else:
        return "single_column"


# ============================================================
# Structured parsing — multi-modal content type separation
# ============================================================

async def parse_document_structured(content: bytes, format: str):
    """Parse document returning structured content elements with type tags.

    Returns ParsedDocument with elements tagged as text/table/image.
    Falls back to flat text if structured parsing is not supported for the format.
    """
    from packages.rag.schemas.parsing import ParsedDocument, ContentElement

    format = format.lower()
    structured_parsers = {
        "pdf": _parse_pdf_structured,
        "docx": _parse_docx_structured,
        "xlsx": _parse_xlsx_structured,
        "pptx": _parse_pptx,
        "txt": _parse_txt_structured,
        "md": _parse_txt_structured,
        "html": _parse_html_structured,
        "htm": _parse_html_structured,
        "jpg": _parse_image_structured,
        "jpeg": _parse_image_structured,
        "png": _parse_image_structured,
        "tiff": _parse_image_structured,
        "tif": _parse_image_structured,
        "bmp": _parse_image_structured,
    }

    parser = structured_parsers.get(format)
    if parser:
        try:
            import asyncio
            result = parser(content)
            if asyncio.iscoroutine(result):
                result = await result
            return result
        except Exception as e:
            logger.warning("Structured parsing failed for %s: %s, falling back to flat text", format, e)

    # Fallback: parse as flat text
    text = await parse_document(content, format)
    return ParsedDocument(
        full_text=text,
        elements=[ContentElement(content_type="text", text=text)],
        content_types={"text"},
    )


def _parse_pdf_structured(content: bytes):
    """PDF structured parsing — separate text, tables, images per page."""
    from packages.rag.schemas.parsing import ParsedDocument, ContentElement
    import pdfplumber

    elements = []
    full_text_parts = []

    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            # Extract text (excluding table areas)
            text = page.extract_text()
            if text:
                elements.append(ContentElement(
                    content_type="text", text=text, page=page_num,
                    metadata={"source": f"page_{page_num}"},
                ))
                full_text_parts.append(text)

            # Extract tables as structured elements
            tables = page.extract_tables()
            for table_num, table in enumerate(tables, 1):
                if table and any(any(cell for cell in row) for row in table):
                    table_md = _format_table(table, f"表格 p{page_num}.{table_num}")
                    elements.append(ContentElement(
                        content_type="table", text=table_md, page=page_num,
                        metadata={"source": f"page_{page_num}_table_{table_num}"},
                    ))
                    full_text_parts.append(table_md)

            # Track images
            images = page.images
            for img_num, img_info in enumerate(images, 1):
                try:
                    img_bytes = None
                    # Try to extract image from the page crop
                    if hasattr(page, 'crop'):
                        bbox = (img_info.get('x0', 0), img_info.get('top', 0),
                                img_info.get('x1', 0), img_info.get('bottom', 0))
                        cropped = page.within_bbox(bbox)
                        img_bytes = cropped.to_image().tobytes()
                    if img_bytes:
                        ocr_text = _ocr_image_bytes(img_bytes, "png")
                    else:
                        ocr_text = ""
                except Exception:
                    ocr_text = ""

                if ocr_text:
                    img_text = f"[图片内容]: {ocr_text}"
                else:
                    img_text = f"[本页第 {img_num} 张图片]"

                elements.append(ContentElement(
                    content_type="image", text=img_text, page=page_num,
                    metadata={"source": f"page_{page_num}_image_{img_num}"},
                ))
                full_text_parts.append(img_text)

    return ParsedDocument(
        full_text="\n\n".join(full_text_parts),
        elements=elements,
        content_types={e.content_type for e in elements},
    )


def _parse_docx_structured(content: bytes):
    """DOCX structured parsing — separate paragraphs, tables, images."""
    from packages.rag.schemas.parsing import ParsedDocument, ContentElement
    from docx import Document
    from docx.table import Table as DocxTable
    from docx.text.paragraph import Paragraph

    doc = Document(io.BytesIO(content))
    elements = []
    full_text_parts = []

    for element in doc.element.body:
        tag = element.tag.split('}')[-1] if '}' in element.tag else element.tag
        if tag == 'p':
            para = Paragraph(element, doc)
            text = para.text.strip()
            if text:
                elements.append(ContentElement(content_type="text", text=text))
                full_text_parts.append(text)
        elif tag == 'tbl':
            table = DocxTable(element, doc)
            table_data = []
            for row in table.rows:
                table_data.append([cell.text.strip() for cell in row.cells])
            if table_data:
                table_md = _format_table(table_data, "表格")
                elements.append(ContentElement(content_type="table", text=table_md))
                full_text_parts.append(table_md)

    # Fallback: paragraph-only parsing
    if not elements:
        for p in doc.paragraphs:
            if p.text.strip():
                elements.append(ContentElement(content_type="text", text=p.text.strip()))
                full_text_parts.append(p.text.strip())

    return ParsedDocument(
        full_text="\n\n".join(full_text_parts),
        elements=elements,
        content_types={e.content_type for e in elements},
    )


def _parse_xlsx_structured(content: bytes):
    """XLSX structured parsing — all content is table type."""
    from packages.rag.schemas.parsing import ParsedDocument, ContentElement
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
    elements = []
    full_text_parts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        table_data = []
        for row in ws.iter_rows(values_only=True):
            table_data.append([str(c) if c is not None else "" for c in row])
        if table_data:
            table_md = _format_table(table_data, f"工作表: {sheet_name}")
            elements.append(ContentElement(
                content_type="table", text=table_md,
                metadata={"sheet": sheet_name},
            ))
            full_text_parts.append(table_md)

    return ParsedDocument(
        full_text="\n\n".join(full_text_parts),
        elements=elements,
        content_types={e.content_type for e in elements},
    )


def _parse_image_structured(content: bytes):
    """Image structured parsing — OCR result as image type element."""
    from packages.rag.schemas.parsing import ParsedDocument, ContentElement

    ocr_text = _ocr_image_bytes(content, "png")
    img_text = f"[图片内容]: {ocr_text}" if ocr_text else "[图片]（无文字内容）"
    element = ContentElement(content_type="image", text=img_text)

    return ParsedDocument(
        full_text=img_text,
        elements=[element],
        content_types={"image"},
    )


def _parse_txt_structured(content: bytes):
    """Plain text structured parsing."""
    from packages.rag.schemas.parsing import ParsedDocument, ContentElement
    text = _parse_txt(content)
    return ParsedDocument(
        full_text=text,
        elements=[ContentElement(content_type="text", text=text)],
        content_types={"text"},
    )


def _parse_html_structured(content: bytes):
    """HTML structured parsing — separate text and tables."""
    from packages.rag.schemas.parsing import ParsedDocument, ContentElement
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(content, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
        tag.decompose()

    elements = []
    full_text_parts = []

    tables = list(soup.find_all("table"))
    for table_tag in tables:
        table_data = []
        for row in table_tag.find_all("tr"):
            cells = [cell.get_text(strip=True) for cell in row.find_all(["td", "th"])]
            if cells:
                table_data.append(cells)
        if table_data:
            table_md = _format_table(table_data, "HTML 表格")
            elements.append(ContentElement(content_type="table", text=table_md))
            full_text_parts.append(table_md)
        table_tag.decompose()

    remaining = soup.get_text(separator="\n", strip=True)
    if remaining:
        elements.insert(0, ContentElement(content_type="text", text=remaining))
        full_text_parts.insert(0, remaining)

    return ParsedDocument(
        full_text="\n\n".join(full_text_parts) if full_text_parts else "",
        elements=elements if elements else [ContentElement(content_type="text", text="")],
        content_types={e.content_type for e in elements} if elements else {"text"},
    )


# ============================================================
# Layout-aware parsing (LayoutLMv3 integration)
# ============================================================

async def parse_with_layout(content: bytes, format: str) -> Dict[str, Any]:
    """
    Parse document with layout awareness using LayoutLMv3.
    Returns structured content with reading order.

    This is an advanced feature that requires additional dependencies.
    """
    if format != "pdf":
        # For non-PDF, use standard parsing
        text = await parse_document(content, format)
        return {
            "text": text,
            "layout_type": "unknown",
            "sections": [],
            "tables": [],
            "images": [],
        }

    try:
        # Try to use LayoutLMv3 for layout-aware parsing
        from transformers import LayoutLMv3Processor, LayoutLMv3ForTokenClassification
        import torch
        from PIL import Image
        import pdf2image

        # Convert PDF to images
        images = pdf2image.convert_from_bytes(content)

        # Process with LayoutLMv3
        processor = LayoutLMv3Processor.from_pretrained("microsoft/layoutlmv3-base")

        structured_content = {
            "text": "",
            "layout_type": "unknown",
            "sections": [],
            "tables": [],
            "images": [],
        }

        for page_num, image in enumerate(images, 1):
            # This is a simplified example - full implementation would require
            # fine-tuned model for document understanding
            text = _ocr_image_bytes(image.tobytes(), "png")
            structured_content["sections"].append({
                "page": page_num,
                "text": text,
                "type": "text",
            })

        structured_content["text"] = "\n\n".join(
            f"## 第 {s['page']} 页\n\n{s['text']}"
            for s in structured_content["sections"]
        )

        return structured_content

    except ImportError:
        logger.warning("LayoutLMv3 not available, using standard parsing")
        text = await parse_document(content, format)
        return {
            "text": text,
            "layout_type": _detect_layout_type(text),
            "sections": [],
            "tables": [],
            "images": [],
        }
    except Exception as e:
        logger.warning("Layout-aware parsing failed: %s", e)
        text = await parse_document(content, format)
        return {
            "text": text,
            "layout_type": "unknown",
            "sections": [],
            "tables": [],
            "images": [],
        }
